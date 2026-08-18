#!/usr/bin/env python3
"""
Build presentation-ready Board deliverables (PPTX, PDF, XLSX, TXT) from the
audited source-of-truth numbers. Read-only: consumes constants below + writes
files into presentation-ready/. No DB access, no secrets.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from fpdf import FPDF
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

OUT = os.path.join(os.path.dirname(__file__), "..", "board-report-2026-06-16", "presentation-ready")
os.makedirs(OUT, exist_ok=True)

# ── SOURCE OF TRUTH (verified live pull 2026-06-06 ~14:38 UTC) ──────────────
TS = "Data: live production pull, 2026-06-06 (UTC). Read-only."
CAPTURED, ROUTED, PRIVATE = 11921, 11020, 901
LOWSTAR, PRIV4, PRIV5 = 606, 211, 84
RESOLVED, ACTIONABLE = 111, 75
D7, D30 = 856, 4222
RATING_RANGE = "4.5-4.8"
ROUTE_RATE = "92.4%"
PRIV_RATE = "7.6%"

# Per-location: name, rating, profile_count, captured, routed, private, lowstar, route_rate, priv_rate, first_review, note, takeaway
LOCS = [
 ["Estancia Angelopolis","4.7",3418,2184,2068,116,112,0.947,0.053,"2026-03-06","Clean; 4 of 116 private are 4-star","Highest volume; working well"],
 ["Estancia Veracruz","4.6",2830,2095,1980,115,115,0.945,0.055,"2026-03-18","Clean; all private are <=3 star","Strong, clean data"],
 ["Estancia Leon","4.7",3185,1772,1577,195,116,0.890,0.110,"2026-03-18","Threshold 5 -> 79 4-star held private","High private rate is threshold, not unhappy guests"],
 ["Harbor's Veracruz","4.6",2225,1197,1150,47,47,0.961,0.039,"2026-03-18","Cleanest location","Best routing rate"],
 ["La Silla Juarez","4.7",2254,1073,979,94,53,0.912,0.088,"2026-03-07","Threshold 5 -> 40 4-star + 1 5-star private","Some positive reviews held private by threshold"],
 ["Harbor's Angelopolis","4.7",2216,906,837,69,35,0.924,0.076,"2026-03-24","Threshold 5 -> 34 4-star held private","Working; private includes 4-star"],
 ["La Silla Huexotitla","4.6",1980,818,752,66,46,0.919,0.081,"2026-03-18","20 4-star private despite threshold 4 (config history)","Minor config history; check routing config"],
 ["Estancia Juarez","4.7",2778,578,554,24,24,0.958,0.042,"2026-03-18","Clean; all private <=3 star","Clean; has a detailed case-study card"],
 ["Estancia Xalapa","4.7",407,436,399,37,12,0.915,0.085,"2026-03-18","FLAG: captured (436) > profile (407); taps != posted reviews; confirm place ID","Looks odd but expected; confirm Google profile"],
 ["Estancia Queretaro","4.6",1806,403,379,24,22,0.940,0.060,"2026-03-16","2 4-star private (minor)","Working normally"],
 ["Regio Norte","4.8",1482,319,213,106,17,0.668,0.332,"2026-03-18","FLAG: 83 5-star private, all March; corrected from April (historical config)","Low routing rate was a March setup issue, now fixed"],
 ["SteakCompany Queretaro","4.5",1659,140,132,8,7,0.943,0.057,"2026-03-18","Lowest volume","Smallest location; working"],
]
assert len(LOCS) == 12, "must be exactly 12 locations"
# math check
for L in LOCS:
    assert L[3] == L[4] + L[5], f"reconcile fail {L[0]}"
assert sum(L[3] for L in LOCS) == CAPTURED
assert sum(L[4] for L in LOCS) == ROUTED
assert sum(L[5] for L in LOCS) == PRIVATE

CLAIMS = [
 ["Live across all 12 Grupo Estancia locations","Yes","12 locations live and capturing","High","Low","Live across all 12 locations"],
 ["11,921 reviews captured","Yes","11,921 (12 locations)","High","Low","Captured 11,921 guest reviews"],
 ["11,020 routed toward Google","Yes","11,020 (12 locations)","High","Low","Routed 11,020 guests toward Google"],
 ["901 held private","Only with caveat","901 = 606 (<=3) + 211 (4) + 84 (5)","High","Medium","901 routed to private feedback; 606 are low (<=3 star)"],
 ["606 low-star private","Yes","606 reviews <=3 star held private","High","Low","606 low ratings routed to a private feedback path"],
 ["111 resolved low-star private","Only with caveat","111 marked resolved (dashboard status)","Medium","High","111 low-star private reviews marked resolved in the dashboard"],
 ["856 reviews in last 7 days","Yes","d7 = 856","High","Low","856 reviews captured in the last 7 days"],
 ["Current Google ratings 4.5-4.8","Yes","range across 12 locations","High","Low","Group Google ratings are 4.5-4.8 stars"],
 ["Google profile counts are RateTap's output","No","Lifetime Google totals predate RateTap","High","Very High","Profile totals are lifetime; RateTap's contribution is 11,020 routed"],
 ["RateTap raised Google ratings","No","Early deltas were a Google place-ID artifact","High","Very High","We do not attribute rating changes to RateTap"],
 ["Bad reviews saved / prevented","No","Cannot prove a guest would have posted","High","High","Low ratings were routed to private feedback instead of the Google link"],
 ["0 active external paying customers","Yes","Only El Braserio had Stripe; now canceled","High","Low","No active external paying customers"],
 ["$0 recurring external revenue","Yes","No active external subscription","High","Low","$0 recurring external revenue today"],
 ["$0 revenue ever collected","No","El Braserio setup-fee not in DB","Low","High","El Braserio setup-fee collection is unverified pending a Stripe check"],
 ["El Braserio pilot did not convert","Yes","created 2026-04-16, trial ended 05-01, canceled","High","Low","One external pilot signed Apr 16 and did not convert"],
 ["16 active accounts","No","16 = 12 locations + 1 owner + 3 regional","High","High","12 live locations; the 16 includes 4 internal admin accounts"],
 ["Activated March 6","Only with caveat","created_at is a synthetic identical import date","High","High","Live since mid-March 2026 (first captured reviews Mar 16-24)"],
 ["2 missing locations from case-studies","No","All 12 in table; only 3 detailed cards","High","Medium","All 12 are listed; 9 lack a detailed card (not 2 missing locations)"],
]

OPEN_Q = [
 ["El Braserio one-time setup fee","Was it actually collected?","Not in DB; Stripe dashboard only","Finance / billing owner"],
 ["True per-location go-live dates","Confirm real launch dates or approve proxy","DB has synthetic 2026-03-06 import date","Data / Ops owner"],
 ["Regio Norte config history","Confirm March issue fixed","Data shows corrected from April","Product / Engineering"],
 ["Estancia Xalapa Google profile","Confirm correct, current place ID","Captured 436 > profile 407","Engineering"],
 ["'2 missing locations' wording","Confirm intent with Memo","All 12 listed; 9 lack detailed cards","Commercialization / Memo"],
]

EXAMPLES = [
 ("Estancia Angelopolis","3-star","\"La comida tardo mucho, algunos platillos en mal termino.\"","food temperature + speed"),
 ("La Silla Juarez","2-star","\"Los banos de discapacitados mas accesibles.\"","accessibility / facilities"),
 ("Harbor's Angelopolis","3-star","\"Se necesita la cerveza mas helada / vasos sin olor.\"","beverage temp + glassware"),
]

# ── colors ──
INK = RGBColor(0x1A,0x1A,0x1A)
MUTE = RGBColor(0x6B,0x6B,0x6B)
ACCENT = RGBColor(0x8B,0x2E,0x2E)   # estancia charcoal-red
GREEN = RGBColor(0x1E,0x6B,0x3A)
LIGHT = RGBColor(0xF4,0xF1,0xEC)
WHITE = RGBColor(0xFF,0xFF,0xFF)

def fmt(n): return f"{n:,}"

# ════════════════════════════════════════════════════════════════════════
# 1. POWERPOINT
# ════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf

def para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, first=False, space_after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    f.name = "Calibri"
    return p

def bg(s, color):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = color

def band(s, color, t, h):
    sh = s.shapes.add_shape(1, Inches(0), Inches(t), SW, Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def caveat(s, text, t=6.75):
    tf = box(s, 0.6, t, 12.1, 0.6)
    para(tf, "Caveat: " + text, 11, MUTE, italic_bold=False, first=True)

# fix caveat helper (no italic kw)
def caveat(s, text, t=6.7):
    tf = box(s, 0.6, t, 12.1, 0.6)
    p = para(tf, "Caveat: " + text, 11, MUTE, first=True)

# Slide 1 — Title
s = slide(); bg(s, LIGHT)
band(s, ACCENT, 0, 0.25)
tf = box(s, 0.8, 2.4, 11.7, 2.5)
para(tf, "RateTap — Board Update", 44, INK, bold=True, first=True)
para(tf, "Grupo Estancia: internal rollout and commercialization decision", 22, MUTE, space_after=18)
para(tf, "Prepared for: Junta de Consejo + Direccion General + CEO  |  Meeting: Tuesday, June 16, 2026", 14, INK)
para(tf, TS, 12, MUTE)
notes(s, "One-line frame: RateTap works internally across all 12 locations; today's decision is whether to keep it internal, run a paid pilot, or commercialize. These are usage figures, not revenue.")

# Slide 2 — Executive summary (5 numbers)
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf = box(s, 0.6, 0.12, 12.1, 0.7); para(tf, "Executive summary", 28, WHITE, bold=True, first=True)
tf = box(s, 0.6, 1.05, 12.1, 0.6)
para(tf, "RateTap is working internally; commercialization is the next decision.", 18, INK, bold=True, first=True)
nums = [("12","live locations"),(fmt(CAPTURED),"reviews captured"),(fmt(ROUTED),"routed toward Google"),(fmt(LOWSTAR),"low ratings routed privately"),("$0","recurring external revenue")]
x = 0.6; cw = 2.35
for val, lab in nums:
    card = s.shapes.add_shape(5, Inches(x), Inches(2.1), Inches(cw), Inches(2.0))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT; card.line.color.rgb = ACCENT; card.line.width=Pt(1)
    card.shadow.inherit=False
    ctf = card.text_frame; ctf.word_wrap=True; ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=ctf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=val
    r.font.size=Pt(30); r.font.bold=True; r.font.color.rgb=ACCENT
    p2=ctf.add_paragraph(); p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run(); r2.text=lab
    r2.font.size=Pt(12); r2.font.color.rgb=INK
    x += cw + 0.18
tf = box(s, 0.6, 4.5, 12.1, 1.8)
para(tf, "Used every day: 856 reviews in the last 7 days; group Google ratings 4.5-4.8 stars.", 15, INK, first=True)
para(tf, "All 12 locations are Grupo Estancia's own. There are no active external paying customers today.", 14, MUTE)
caveat(s, "These are usage and routing figures, not revenue, and not a claim of Google-rating growth.")
notes(s, "Lead with what's proven, then state pre-revenue immediately so no one infers revenue. The five numbers are the whole story in one glance.")

# Slide 3 — What RateTap does (flow)
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf = box(s, 0.6, 0.12, 12.1, 0.7); para(tf, "What RateTap does", 28, WHITE, bold=True, first=True)
steps = ["Guest taps card","Leaves a star rating","Satisfied guest\nrouted to Google","Lower rating\nrouted privately","Manager follows up"]
x=0.55; cw=2.3
for i,st in enumerate(steps):
    col = GREEN if "Google" in st else (ACCENT if "privately" in st else INK)
    card=s.shapes.add_shape(5, Inches(x), Inches(2.5), Inches(cw), Inches(1.6))
    card.fill.solid(); card.fill.fore_color.rgb=LIGHT; card.line.color.rgb=col; card.line.width=Pt(1.5); card.shadow.inherit=False
    ctf=card.text_frame; ctf.word_wrap=True; ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=ctf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=st
    r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=col
    if i < len(steps)-1:
        ar=s.shapes.add_shape(1, Inches(x+cw-0.03), Inches(3.27), Inches(0.16), Inches(0.06))
        ar.fill.solid(); ar.fill.fore_color.rgb=MUTE; ar.line.fill.background(); ar.shadow.inherit=False
    x+=cw+0.0
tf=box(s,0.6,4.6,12.1,1.4)
para(tf,"The threshold is set per location. Six locations use a 5-star threshold, so some 4-star guests are routed to the private form.",14,INK,first=True)
caveat(s,"\"Routed privately\" is a routing rule; it does not always mean the guest was unhappy.")
notes(s,"This mechanic underpins every later number. Happy -> Google, lower -> private form for the manager.")

# Slide 4 — Current rollout
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"Current rollout — all 12 locations live",26,WHITE,bold=True,first=True)
brands={}
for L in LOCS:
    b=L[0].split()[0] if not L[0].startswith("La Silla") and not L[0].startswith("Harbor") else (" ".join(L[0].split()[:2]))
for grp,members in [("Estancia (6)",["Angelopolis","Juarez","Queretaro","Leon","Veracruz","Xalapa"]),
                    ("Harbor's (2)",["Angelopolis","Veracruz"]),
                    ("La Silla (2)",["Juarez","Huexotitla"]),
                    ("Others (2)",["SteakCompany Queretaro","Regio Norte"])]:
    pass
tf=box(s,0.6,1.15,5.9,5.2)
para(tf,"Brands & locations",16,ACCENT,bold=True,first=True)
para(tf,"Estancia (6): Angelopolis, Juarez, Queretaro, Leon, Veracruz, Xalapa",13,INK,space_after=8)
para(tf,"Harbor's (2): Angelopolis, Veracruz",13,INK,space_after=8)
para(tf,"La Silla (2): Juarez, Huexotitla",13,INK,space_after=8)
para(tf,"SteakCompany (1): Queretaro",13,INK,space_after=8)
para(tf,"Regio Norte (1)",13,INK,space_after=8)
tf=box(s,6.8,1.15,5.9,5.2)
para(tf,"Status",16,ACCENT,bold=True,first=True)
para(tf,"12 locations live and capturing daily",14,INK,space_after=8)
para(tf,"Group Google ratings: 4.5-4.8 stars",14,INK,space_after=8)
para(tf,"Live since mid-March 2026",14,INK,space_after=8)
para(tf,"Note: the database shows \"16 active\" because it includes 4 internal admin accounts (1 owner + 3 regional). The real live-location count is 12.",13,MUTE,space_after=8)
caveat(s,"Google ratings shown are current profile ratings, not RateTap-caused growth.")
notes(s,"Clarify the 16-vs-12 point proactively. All 12 are our own restaurants.")

# Slide 5 — Capture and routing (bar)
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"Review capture and routing",28,WHITE,bold=True,first=True)
# simple horizontal bars
rows=[("Captured",CAPTURED,INK),("Routed toward Google",ROUTED,GREEN),("Held private",PRIVATE,ACCENT),("Low rating (<=3*) private",LOWSTAR,ACCENT)]
maxv=CAPTURED; y=1.5; barmax=8.3
for lab,val,col in rows:
    tf=box(s,0.6,y-0.05,3.3,0.5); para(tf,lab,13,INK,bold=True,first=True)
    w=max(0.15, barmax*val/maxv)
    bar=s.shapes.add_shape(5, Inches(4.0), Inches(y), Inches(w), Inches(0.5))
    bar.fill.solid(); bar.fill.fore_color.rgb=col; bar.line.fill.background(); bar.shadow.inherit=False
    tf=box(s,4.0+w+0.1,y-0.05,2.0,0.5); para(tf,fmt(val),13,INK,bold=True,first=True)
    y+=0.95
tf=box(s,0.6,5.5,12.1,1.0)
para(tf,f"Routing rate {ROUTE_RATE} of captured guests. Held private {PRIV_RATE}. (901 private = 606 low + 211 four-star + 84 five-star.)",13,MUTE,first=True)
caveat(s,"Routed toward Google does not guarantee the guest completed a public post.")
notes(s,"92.4% of captured guests were routed to Google. State that 'held private' is not all bad reviews.")

# Slide 6 — Private feedback loop
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"Private feedback loop",28,WHITE,bold=True,first=True)
stt=[(fmt(LOWSTAR),"low ratings (<=3*)\nrouted privately"),(fmt(RESOLVED),"marked resolved\nin the dashboard"),("~"+str(ACTIONABLE),"with substantive\nwritten feedback")]
x=0.6;cw=3.0
for v,l in stt:
    card=s.shapes.add_shape(5,Inches(x),Inches(1.3),Inches(cw),Inches(1.5))
    card.fill.solid();card.fill.fore_color.rgb=LIGHT;card.line.color.rgb=ACCENT;card.line.width=Pt(1);card.shadow.inherit=False
    ctf=card.text_frame;ctf.word_wrap=True;ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=ctf.paragraphs[0];p.alignment=PP_ALIGN.CENTER;r=p.add_run();r.text=v;r.font.size=Pt(26);r.font.bold=True;r.font.color.rgb=ACCENT
    p2=ctf.add_paragraph();p2.alignment=PP_ALIGN.CENTER;r2=p2.add_run();r2.text=l;r2.font.size=Pt(11);r2.font.color.rgb=INK
    x+=cw+0.2
tf=box(s,0.6,3.2,12.1,3.0)
para(tf,"Examples of actionable private feedback (routed to managers, not posted publicly):",14,ACCENT,bold=True,first=True)
for loc,star,quote,cat in EXAMPLES:
    para(tf,f"{loc} ({star}) — {quote}  [{cat}]",13,INK,space_after=8)
caveat(s,"\"Resolved\" is a dashboard status; it does not prove a public bad review was prevented.")
notes(s,"Say 'routed privately', never 'saved' or 'prevented'. Many low-star taps are silent or mis-taps with positive text.")

# Slide 7 — Scorecard table
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"Location scorecard",28,WHITE,bold=True,first=True)
cols=["Location","Google*","Captured","Routed","Private","Note"]
rows=1+len(LOCS)+1
tbl=s.shapes.add_table(rows,len(cols),Inches(0.4),Inches(1.05),Inches(12.5),Inches(6.0)).table
widths=[3.1,1.0,1.3,1.3,1.1,4.7]
for i,w in enumerate(widths): tbl.columns[i].width=Inches(w)
for j,c in enumerate(cols):
    cell=tbl.cell(0,j); cell.text=c
    cell.fill.solid(); cell.fill.fore_color.rgb=ACCENT
    pr=cell.text_frame.paragraphs[0]; pr.runs[0].font.size=Pt(11); pr.runs[0].font.bold=True; pr.runs[0].font.color.rgb=WHITE
for i,L in enumerate(LOCS,start=1):
    vals=[L[0],L[1],fmt(L[3]),fmt(L[4]),fmt(L[5]), ("FLAG: "+L[10].split(';')[0]) if L[10].startswith("FLAG") else L[10].split(';')[0]]
    for j,v in enumerate(vals):
        cell=tbl.cell(i,j); cell.text=str(v)
        run=cell.text_frame.paragraphs[0].runs[0]; run.font.size=Pt(9.5)
        run.font.color.rgb = ACCENT if (j==5 and str(v).startswith("FLAG")) else INK
        if i%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=LIGHT
        else: cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
tot=tbl.cell(rows-1,0); tot.text="TOTAL (12)"
totvals=["", fmt(CAPTURED),fmt(ROUTED),fmt(PRIVATE),"captured = routed + private"]
for j,v in enumerate(totvals,start=1):
    cell=tbl.cell(rows-1,j); cell.text=str(v)
for j in range(len(cols)):
    cell=tbl.cell(rows-1,j); cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xE6,0xDD,0xCE)
    if cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].font.size=Pt(10); cell.text_frame.paragraphs[0].runs[0].font.bold=True
tf=box(s,0.4,7.05,12.5,0.4); para(tf,"*Current Google profile rating (not RateTap-caused). Full appendix in the Excel workbook.",9,MUTE,first=True)
notes(s,"Two locations carry flags: Regio Norte (March config, now fixed) and Xalapa (captured > profile; confirm place ID). Everything reconciles.")

# Slide 8 — Proven vs not proven
s = slide(); bg(s, WHITE)
band(s, GREEN, 0, 0.0)
tf=box(s,0.6,0.3,12.1,0.7); para(tf,"What is proven vs. not proven",28,INK,bold=True,first=True)
# proven panel
p1=s.shapes.add_shape(5,Inches(0.6),Inches(1.3),Inches(5.9),Inches(5.0))
p1.fill.solid();p1.fill.fore_color.rgb=RGBColor(0xEA,0xF3,0xEC);p1.line.color.rgb=GREEN;p1.line.width=Pt(1.5);p1.shadow.inherit=False
ctf=p1.text_frame;ctf.word_wrap=True;ctf.margin_left=Inches(0.25);ctf.margin_top=Inches(0.2)
para(ctf,"PROVEN",18,GREEN,bold=True,first=True)
for t in ["Internal usage across 12 locations","Active daily usage (856 in 7 days)","Thousands routed toward Google (11,020)","A working private feedback loop"]:
    para(ctf,"+  "+t,14,INK,space_after=10)
p2=s.shapes.add_shape(5,Inches(6.8),Inches(1.3),Inches(5.9),Inches(5.0))
p2.fill.solid();p2.fill.fore_color.rgb=RGBColor(0xF7,0xEC,0xEC);p2.line.color.rgb=ACCENT;p2.line.width=Pt(1.5);p2.shadow.inherit=False
ctf=p2.text_frame;ctf.word_wrap=True;ctf.margin_left=Inches(0.25);ctf.margin_top=Inches(0.2)
para(ctf,"NOT PROVEN",18,ACCENT,bold=True,first=True)
for t in ["External SaaS demand","Google-rating growth caused by RateTap","Public bad reviews prevented","External recurring revenue ($0 today)"]:
    para(ctf,"-  "+t,14,INK,space_after=10)
notes(s,"This is the honest heart of the deck. Operational value is real; external business is unproven.")

# Slide 9 — Board decision
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"Board decision",28,WHITE,bold=True,first=True)
opts=[("Option A","Keep internal only","No cost, no risk. No revenue line.",LIGHT,INK),
      ("Option B  (Recommended)","Run 1-2 controlled paid pilots","Lowest-risk way to test willingness to pay. Builds a reference.",RGBColor(0xEA,0xF3,0xEC),GREEN),
      ("Option C","Invest in full commercialization","Biggest upside, but premature. Demand unproven.",LIGHT,INK)]
y=1.3
for title,sub,desc,fillc,col in opts:
    card=s.shapes.add_shape(5,Inches(0.6),Inches(y),Inches(12.1),Inches(1.5))
    card.fill.solid();card.fill.fore_color.rgb=fillc;card.line.color.rgb=col;card.line.width=Pt(1.5 if col==GREEN else 0.75);card.shadow.inherit=False
    ctf=card.text_frame;ctf.word_wrap=True;ctf.margin_left=Inches(0.3);ctf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=ctf.paragraphs[0];r=p.add_run();r.text=f"{title}:  {sub}";r.font.size=Pt(17);r.font.bold=True;r.font.color.rgb=col
    p2=ctf.add_paragraph();r2=p2.add_run();r2.text=desc;r2.font.size=Pt(12);r2.font.color.rgb=INK
    y+=1.7
tf=box(s,0.6,6.55,12.1,0.7)
para(tf,"Recommendation: Option B — validate willingness to pay before over-investing; keeps the product focused.",13,ACCENT,bold=True,first=True)
notes(s,"Recommend B. Pre-revenue is the next phase, not failure. Ask for the go/no-go and a pilot mandate.")

# Slide 10 — 30/60/90
s = slide(); bg(s, WHITE)
band(s, ACCENT, 0, 0.9)
tf=box(s,0.6,0.12,12.1,0.7); para(tf,"30 / 60 / 90-day plan",28,WHITE,bold=True,first=True)
plan=[("First 30 days",["Confirm open facts (incl. billing)","Onboarding + config QA checklist","Select pilot targets"]),
      ("By 60 days",["Launch 1-2 paid pilots","Validate billing collection","Monitor usage and retention"]),
      ("By 90 days",["Decide whether to scale","Build a referenceable case study","Define pricing & support model"])]
x=0.6;cw=3.95
for title,items in plan:
    card=s.shapes.add_shape(5,Inches(x),Inches(1.3),Inches(cw),Inches(4.6))
    card.fill.solid();card.fill.fore_color.rgb=LIGHT;card.line.color.rgb=ACCENT;card.line.width=Pt(1);card.shadow.inherit=False
    ctf=card.text_frame;ctf.word_wrap=True;ctf.margin_left=Inches(0.25);ctf.margin_top=Inches(0.25)
    para(ctf,title,18,ACCENT,bold=True,first=True)
    for it in items: para(ctf,"-  "+it,14,INK,space_after=12)
    x+=cw+0.2
caveat(s,"Open facts to confirm first include El Braserio billing (Stripe) and true go-live dates.")
notes(s,"Concrete, sober, reversible. The 90-day gate decides Option C.")

pptx_path=os.path.join(OUT,"RateTap_Board_Deck_2026-06-16.pptx")
prs.save(pptx_path)
print("PPTX:", pptx_path, len(prs.slides._sldIdLst), "slides")

# ════════════════════════════════════════════════════════════════════════
# 2. PDFs (board packet, speaker notes, readme) via fpdf2
# ════════════════════════════════════════════════════════════════════════
def clean(t):
    # strip chars latin-1 can't encode
    return (t.replace("→","->").replace("≤","<=").replace("—","-")
             .replace("–","-").replace("’","'").replace("“",'"')
             .replace("”",'"').replace("‘","'").replace("★","*")
             .replace("…","...").replace("á","a").replace("é","e")
             .replace("í","i").replace("ó","o").replace("ú","u")
             .replace("ñ","n").replace("Á","A").replace("ç","c")
             .replace("≈","~").replace("️","").replace("\U0001f44c",""))

class PDF(FPDF):
    title_txt=""
    def header(self):
        if self.page_no()==1: return
        self.set_font("Helvetica","I",8); self.set_text_color(140,140,140)
        self.set_x(self.l_margin)
        self.cell(95,8,clean(self.title_txt),align="L")
        self.cell(0,8,"RateTap - Confidential",align="R",new_x="LMARGIN",new_y="NEXT")
        self.set_text_color(26,26,26); self.ln(2)
    def footer(self):
        if self.page_no()==1: return
        self.set_y(-12); self.set_font("Helvetica","I",8); self.set_text_color(140,140,140)
        self.set_x(self.l_margin)
        self.cell(0,8,f"Page {self.page_no()-1}",align="C")

def mc(pdf,h,t,size,style="",fill=False,indent=0):
    pdf.set_font("Helvetica",style,size); pdf.set_x(pdf.l_margin+indent)
    pdf.multi_cell(0,h,clean(t),border=0,align="L",fill=fill,new_x="LMARGIN",new_y="NEXT")
def h1(pdf,t):
    pdf.set_text_color(139,46,46); pdf.ln(2); mc(pdf,8,t,15,"B"); pdf.set_text_color(26,26,26); pdf.ln(1)
def h2(pdf,t):
    pdf.set_text_color(40,40,40); pdf.ln(1.5); mc(pdf,6,t,11.5,"B"); pdf.set_text_color(26,26,26)
def body(pdf,t):
    pdf.set_text_color(26,26,26); mc(pdf,5.4,t,10.5)
def bullet(pdf,t):
    mc(pdf,5.4,"- "+t,10.5,indent=4)
def boxed(pdf,title,lines,fill=(244,241,236)):
    pdf.ln(1); pdf.set_fill_color(*fill); pdf.set_draw_color(139,46,46)
    mc(pdf,6,title,10.5,"B",fill=True)
    for ln in lines: mc(pdf,5,"  "+ln,9.8,fill=True)
    pdf.ln(1)

# ---- 2a. Board packet PDF ----
pdf=PDF(); pdf.set_compression(False); pdf.title_txt="RateTap Board / DG Report"; pdf.set_auto_page_break(True,15)
def ctr(p,h,txt,size,style=""):
    p.set_font("Helvetica",style,size); p.set_x(p.l_margin)
    p.multi_cell(0,h,clean(txt),border=0,align="C",new_x="LMARGIN",new_y="NEXT")
pdf.add_page()  # cover
pdf.set_fill_color(139,46,46); pdf.rect(0,0,210,6,"F")
pdf.ln(60); pdf.set_text_color(26,26,26)
ctr(pdf,14,"RateTap",30,"B")
ctr(pdf,10,"Board / DG / CEO Report",18,"B")
pdf.ln(6); pdf.set_text_color(90,90,90)
ctr(pdf,8,"Grupo Estancia: internal rollout and commercialization decision",13)
pdf.ln(20)
ctr(pdf,7,"Meeting: Tuesday, June 16, 2026",11)
ctr(pdf,7,TS,11)
ctr(pdf,7,"Confidential - internal",11)

pdf.add_page()
h1(pdf,"1. Executive summary")
for t in ["RateTap is live across all 12 Grupo Estancia locations and used every day.",
          f"It has captured {fmt(CAPTURED)} reviews, routed {fmt(ROUTED)} toward Google, and routed {fmt(LOWSTAR)} low ratings privately.",
          "Group Google ratings are 4.5-4.8 stars; we do not claim RateTap raised them.",
          "It is pre-revenue externally: no active external paying customers.",
          "Board decision: keep internal, run a paid pilot (recommended), or invest in full commercialization."]:
    bullet(pdf,t)
boxed(pdf,"Definitions",
      ["Captured = a guest tapped the card and rated.",
       "Routed to Google = sent to the Google review page.",
       "Held private = sent to a private feedback form instead.",
       "Low-star private = a <=3-star rating held private.",
       "Google profile count = Google's lifetime total (not RateTap's)."])
boxed(pdf,"Verified vs. not verified",
      [f"VERIFIED: 12 live locations; {fmt(CAPTURED)} captured; {fmt(ROUTED)} routed; {fmt(PRIVATE)} held private; {fmt(LOWSTAR)} low-star private; {RESOLVED} resolved; {D7} in 7 days; ratings 4.5-4.8; 0 active external paying customers; $0 recurring external revenue.",
       "NOT VERIFIED (needs confirmation): whether El Braserio's setup fee was collected; true go-live dates; Xalapa's Google place ID."],
      fill=(234,243,236))

h1(pdf,"2. What RateTap is")
body(pdf,"A simple feedback system. A guest taps a card at the table and rates the visit. Happy guests are sent to Google; lower-rating guests are sent to a private form that reaches a manager.")
h1(pdf,"3. Why Grupo Estancia built it")
for t in ["Capture happy guests while they are still in the restaurant.","Route satisfied guests toward Google reviews.","Route lower-rating experiences privately so managers can respond.","Give operators a feedback loop they did not have before."]:
    bullet(pdf,t)
h1(pdf,"4. How it works today")
body(pdf,"Guest taps card -> leaves a star rating -> high rating routed to Google; lower rating enters a private feedback form -> manager can review and resolve.")

h1(pdf,"5. Current implementation")
for t in ["Live across all 12 Grupo Estancia locations.",
          "The database shows '16 active', but that includes 4 internal accounts (1 owner + 3 regional). The real live-location count is 12.",
          "This is not an external customer count. All 12 are Grupo Estancia's own restaurants.",
          "Config caveat: 6 locations use a 5-star threshold, so some 4-star guests go to private feedback. Regio Norte had a March setup issue (corrected from April)."]:
    bullet(pdf,t)

h1(pdf,"6. Results to date (verified)")
res=[("Captured by RateTap",fmt(CAPTURED)),("Routed toward Google",f"{fmt(ROUTED)} ({ROUTE_RATE})"),
     ("Held private",fmt(PRIVATE)),("Low-star (<=3) private",fmt(LOWSTAR)),
     ("Resolved low-star private",str(RESOLVED)),("Reviews last 7 days",str(D7))]
pdf.set_font("Helvetica","",10.5)
for k,v in res:
    pdf.set_x(pdf.l_margin)
    pdf.cell(90,6,clean("  "+k),border="B")
    pdf.cell(95,6,clean(v),border="B",new_x="LMARGIN",new_y="NEXT")
pdf.ln(1)

h1(pdf,"7. What the numbers mean")
body(pdf,"The system is genuinely working and in daily use. Most guests are happy and routed to Google. Lower-rating guests are caught in a private channel managers can act on.")
h1(pdf,"8. What the numbers do NOT mean")
for t in ["The Google ratings shown are current profile ratings, not RateTap-caused growth.",
          "Google review counts are lifetime profile totals, not RateTap's output.",
          "Held private does not always mean negative (211 are 4-star, 84 are 5-star).",
          "Resolved does not prove a public bad review was prevented.",
          "There are no active external paying customers today."]:
    bullet(pdf,t)

h1(pdf,"9. Operational value")
body(pdf,f"About {ACTIONABLE} low-star private reviews carry substantive written feedback managers can act on. Examples (routed to managers, not posted publicly):")
for loc,star,quote,cat in EXAMPLES:
    bullet(pdf,f"{loc} ({star}): {quote} [{cat}]")
body(pdf,"We do not claim these would otherwise have been posted publicly.")

h1(pdf,"10. Commercialization status")
for t in ["Internal usage: proven.","External: pre-revenue. One external pilot (El Braserio) signed up in April and did not convert.",
          "A paid pilot is needed to validate that outside groups will pay.",
          "Whether El Braserio's setup fee was collected must be confirmed in Stripe (not in our data)."]:
    bullet(pdf,t)

h1(pdf,"11. Board decision")
for title,desc in [("Option A - Keep internal","No cost, no risk; no revenue line."),
                   ("Option B - Run 1-2 controlled paid pilots (RECOMMENDED)","Tests willingness to pay cheaply; builds a reference; reversible."),
                   ("Option C - Invest in full commercialization","Biggest upside but premature; demand unproven.")]:
    h2(pdf,title); body(pdf,desc)
body(pdf,"Recommendation: Option B. External willingness to pay is unproven (one churned pilot, no confirmed revenue); a paid pilot is the lowest-risk way to test it.")

h1(pdf,"12. Risks and caveats")
for t in ["Usage metrics, not revenue.","No Google-rating-growth claims (place-ID artifact).","Held private includes positive reviews.","Resolved is a status, not proof.","Activation dates are synthetic; we use first captured review.","El Braserio fee collection unverified."]:
    bullet(pdf,t)

h1(pdf,"13. 30 / 60 / 90-day plan")
h2(pdf,"First 30 days");
for t in ["Confirm open facts (incl. El Braserio billing)","Onboarding + config QA checklist","Select pilot targets"]: bullet(pdf,t)
h2(pdf,"By 60 days")
for t in ["Launch 1-2 paid pilots","Validate billing collection","Monitor usage and retention"]: bullet(pdf,t)
h2(pdf,"By 90 days")
for t in ["Decide whether to scale commercialization","Build a referenceable external case study","Define pricing and support model"]: bullet(pdf,t)

h1(pdf,"14. Appendix - location scorecard")
pdf.set_font("Helvetica","B",8)
heads=["Location","Rating","Captured","Routed","Private","Route%"]
w=[52,16,24,24,22,20]
for i,hh in enumerate(heads): pdf.cell(w[i],6,clean(hh),border=1,align="C")
pdf.ln()
pdf.set_font("Helvetica","",8)
for L in LOCS:
    pdf.cell(w[0],5.5,clean(L[0]),border=1)
    pdf.cell(w[1],5.5,L[1],border=1,align="C")
    pdf.cell(w[2],5.5,fmt(L[3]),border=1,align="R")
    pdf.cell(w[3],5.5,fmt(L[4]),border=1,align="R")
    pdf.cell(w[4],5.5,fmt(L[5]),border=1,align="R")
    pdf.cell(w[5],5.5,f"{L[7]*100:.1f}%",border=1,align="R"); pdf.ln()
pdf.set_font("Helvetica","B",8)
pdf.cell(w[0],5.5,"TOTAL (12)",border=1)
pdf.cell(w[1],5.5,"4.5-4.8",border=1,align="C")
pdf.cell(w[2],5.5,fmt(CAPTURED),border=1,align="R")
pdf.cell(w[3],5.5,fmt(ROUTED),border=1,align="R")
pdf.cell(w[4],5.5,fmt(PRIVATE),border=1,align="R")
pdf.cell(w[5],5.5,ROUTE_RATE,border=1,align="R"); pdf.ln()

h1(pdf,"15. Q&A (short)")
qa=[("How much revenue?","$0 recurring external revenue. One pilot did not convert; its setup-fee collection needs a Stripe check."),
    ("Are these paying customers?","No. All 12 are our own locations."),
    ("Did it raise Google ratings?","We do not claim that; early growth figures were a Google data artifact."),
    ("How many reviews did RateTap generate?","It routed 11,020 guests to Google; profile totals are lifetime and predate RateTap."),
    ("Proof bad reviews were prevented?","We routed 606 low ratings to a private form; we cannot prove each guest would have posted.")]
for q,a in qa: h2(pdf,"Q: "+q); body(pdf,"A: "+a)

h1(pdf,"16. Language guardrails")
h2(pdf,"Do NOT say")
for t in ["16 active customers/accounts","RateTap increased/grew our Google ratings","We saved/prevented/intercepted bad reviews","We have paying customers / $X revenue","$0 has ever been collected","Activated March 6"]: bullet(pdf,t)
h2(pdf,"Safe to say")
for t in ["Live across all 12 locations","Captured 11,921; routed 11,020; 606 low ratings routed privately","Ratings 4.5-4.8; we capture and route","No active external paying customers; one pilot churned"]: bullet(pdf,t)

h1(pdf,"17. Open confirmations before June 16")
for q in OPEN_Q: bullet(pdf,f"{q[0]}: {q[1]} ({q[2]})")

packet_pdf=os.path.join(OUT,"RateTap_Board_Packet_2026-06-16.pdf"); pdf.output(packet_pdf)
print("PDF packet:", packet_pdf, pdf.page_no(),"pages")

# ---- 2b. Speaker notes PDF ----
sp=PDF(); sp.set_compression(False); sp.title_txt="RateTap Speaker Notes"; sp.set_auto_page_break(True,15)
sp.add_page(); sp.ln(50); sp.set_text_color(26,26,26)
ctr(sp,12,"RateTap - Speaker Notes",26,"B")
sp.set_text_color(90,90,90); sp.ln(4)
ctr(sp,7,"For Board / DG / CEO meeting - June 16, 2026",12)
ctr(sp,7,TS,12); sp.set_text_color(26,26,26)
sp.add_page()
SCRIPTS=[("15-second version",
  f"RateTap is our in-house review system. It's live in all 12 locations, used daily, and has routed {fmt(ROUTED)} guests to Google. It works internally, but it isn't an external business yet - no paying customers. The decision is whether to run a paid pilot."),
 ("30-second version",
  f"RateTap is the review system we built in-house. It's live across all 12 locations and used every day - almost 12,000 reviews captured, {fmt(ROUTED)} guests routed to Google, and {fmt(LOWSTAR)} low ratings routed to a private form for managers. Our ratings are 4.5-4.8. It's not generating revenue yet - these are our own restaurants - but it's proven internally. The question for the Board is whether we commercialize it."),
 ("90-second version",
  "Three months ago we rolled out RateTap to all 12 locations. A guest taps a card and rates the visit; happy guests go to Google, lower ratings go to a private form that reaches the manager. We've captured 11,921 reviews - about 4,200 a month, 856 in the last week. 11,020 guests routed to Google; 901 to private feedback, of which 606 were genuinely low and about 75 carry actionable comments. Two honest caveats: we do not claim RateTap raised our ratings - the early growth numbers were a Google data artifact; and it's pre-revenue, with no external paying customers and one external pilot that didn't convert. The product works and we have a real case study. The decision is whether to fund commercialization, starting with a controlled paid pilot."),
 ("3-minute version",
  "I'll give a clear, honest picture and the one decision we need. WHAT IT IS: a feedback system we built in-house. A guest taps a card and rates; at or above a location's threshold they go to Google, below it they go to a private form for the manager. WHY: capture happy guests in the moment, route complaints privately. WHERE WE ARE: live across all 12 locations, daily use - 856 reviews in seven days. Since mid-March we captured 11,921; routed 11,020 to Google; 901 to private feedback. THREE PRECISE POINTS: those 901 are not all bad reviews - only 606 are three stars or below; the big Google counts are lifetime totals that predate this system - our contribution is the 11,020 routed; and we do not claim RateTap raised ratings. OPERATIONAL VALUE: about 75 low-star private reviews carry specific feedback a manager can act on. COMMERCIALLY: pre-revenue, all 12 are our own, no external paying customers, one pilot did not convert and its fee collection needs a Stripe check. THE DECISION: keep internal, run a paid pilot, or invest fully. We recommend the paid pilot - lowest-risk way to test willingness to pay before committing capital."),
 ("Spanish-ready / bilingual (90s)",
  "Hace tres meses lanzamos RateTap en las 12 sucursales. El cliente toca una tarjeta y califica; los satisfechos van a Google, las calificaciones mas bajas van a un formulario privado que llega al gerente. Hemos captado 11,921 resenas - unas 4,200 al mes, 856 en la ultima semana. 11,020 dirigidos a Google; 901 a retroalimentacion privada, de las cuales 606 eran bajas y unas 75 traen comentarios accionables. Dos aclaraciones: no afirmamos que RateTap subio las calificaciones - los aumentos iniciales son un artefacto de place-ID de Google; y esta en etapa previa a ingresos, sin clientes externos de pago, y el unico piloto externo no se convirtio. El producto funciona y tenemos un caso real. La decision es si lo comercializamos, empezando por un piloto pagado controlado.")]
for t,txt in SCRIPTS:
    h1(sp,t); body(sp,txt)
h1(sp,'"If challenged" answers')
ifc=[("Is this making money?","No recurring external revenue today. One pilot did not convert; its setup-fee collection needs a Stripe check."),
     ("Are these external customers?","No. All 12 are our own locations."),
     ("Did this increase Google ratings?","We do not claim that. Early growth figures were a Google data artifact."),
     ("How many Google reviews did RateTap generate?","It routed 11,020 guests; profile totals are lifetime and predate RateTap."),
     ("Did we prevent bad Google reviews?","We routed 606 low ratings to a private form; we cannot prove each guest would have posted."),
     ("Why are some 5-star/positive reviews held private?","Threshold settings (6 locations use 5-star) plus occasional mis-taps."),
     ("What is the ask?","A decision: keep internal, run a paid pilot (recommended), or commercialize."),
     ("What's needed before selling externally?","Validated billing, standardized onboarding, config QA, pilot retention data, a reference customer, clear pricing.")]
for q,a in ifc: h2(sp,"Q: "+q); body(sp,"A: "+a)
h1(sp,"Five things NOT to say")
for t in ["'saved / prevented / intercepted' reviews","'increased / grew our Google rating'","'paying customers' or 'revenue of $X'","'16 customers/accounts'","'activated March 6'"]: bullet(sp,t)
h1(sp,"Five open confirmations before June 16")
for q in OPEN_Q: bullet(sp,f"{q[0]} - {q[1]}")
notes_pdf=os.path.join(OUT,"RateTap_Speaker_Notes_2026-06-16.pdf"); sp.output(notes_pdf)
print("PDF notes:", notes_pdf, sp.page_no(),"pages")

# ---- 2c. README one-pager ----
rd=PDF(); rd.set_compression(False); rd.title_txt="Open This First"; rd.set_auto_page_break(True,12)
rd.add_page()
rd.set_fill_color(139,46,46); rd.rect(0,0,210,6,"F"); rd.ln(6)
rd.set_text_color(26,26,26); mc(rd,9,"RateTap Board Brief - Open This First",20,"B")
rd.set_text_color(120,120,120); mc(rd,5,TS,9,"I"); rd.set_text_color(26,26,26); rd.ln(1)
h2(rd,"Safest 5-sentence story")
body(rd,clean("RateTap is an internal review-capture and routing system built for Grupo Estancia and currently live across all 12 of its locations. As of the latest production pull, it has captured 11,921 guest reviews, routed 11,020 guests toward Google, and routed 606 low-rating experiences (<=3 stars) to a private feedback path. It is used daily (856 reviews in the last 7 days) and group Google ratings are 4.5-4.8 stars - though we do not claim RateTap raised them. It is operationally proven but pre-revenue as an outside product: no active external paying customers and $0 recurring external revenue today. The Board decision is whether to keep it internal, run a controlled paid pilot (recommended), or invest in full commercialization."))
h2(rd,"Five numbers that matter")
for t in ["12 live Grupo Estancia locations",f"{fmt(CAPTURED)} guest reviews captured",f"{fmt(ROUTED)} guests routed toward Google",f"{fmt(LOWSTAR)} low-rating experiences routed privately","$0 recurring external revenue today"]: bullet(rd,t)
h2(rd,"What is proven")
for t in ["Live across 12 locations, used daily","Thousands routed toward Google","A working private feedback loop"]: bullet(rd,t)
h2(rd,"What is NOT proven")
for t in ["Google-rating growth caused by RateTap","Public bad reviews prevented","External SaaS demand / revenue"]: bullet(rd,t)
h2(rd,"Board decision requested")
body(rd,"A) Keep internal   B) Run 1-2 paid pilots (RECOMMENDED)   C) Full commercialization")
h2(rd,"Do not say")
body(rd,"'saved/prevented reviews'; 'increased Google rating'; 'paying customers/revenue'; '16 customers'; 'activated March 6'.")
h2(rd,"Needs confirmation before June 16")
for q in OPEN_Q: bullet(rd,f"{q[0]}: {q[1]}")
readme_pdf=os.path.join(OUT,"README_OPEN_THIS_FIRST.pdf"); rd.output(readme_pdf)
print("PDF readme:", readme_pdf, rd.page_no(),"pages")

# ════════════════════════════════════════════════════════════════════════
# 3. EXCEL workbook
# ════════════════════════════════════════════════════════════════════════
wb=openpyxl.Workbook()
hdr_fill=PatternFill("solid",fgColor="8B2E2E"); hdr_font=Font(bold=True,color="FFFFFF")
tot_fill=PatternFill("solid",fgColor="E6DDCE"); thin=Side(style="thin",color="CCCCCC")
border=Border(left=thin,right=thin,top=thin,bottom=thin)
def style_header(ws,ncol):
    for c in range(1,ncol+1):
        cell=ws.cell(1,c); cell.fill=hdr_fill; cell.font=hdr_font; cell.alignment=Alignment(horizontal="center",vertical="center",wrap_text=True); cell.border=border
    ws.freeze_panes="A2"
def autofit(ws):
    for col in ws.columns:
        m=0; letter=get_column_letter(col[0].column)
        for cell in col:
            if cell.value is not None: m=max(m,len(str(cell.value)))
        ws.column_dimensions[letter].width=min(max(m+2,10),52)

# Tab A
ws=wb.active; ws.title="Location Scorecard"
cols=["Location","Google rating","Google profile review count","Captured by RateTap","Routed to Google","Held private","Low-star private","Routing rate","Private rate","First captured review / proxy","Date caveat","Status","Notes / anomalies","Plain-English takeaway"]
ws.append(cols)
for L in LOCS:
    ws.append([L[0],float(L[1]),L[2],L[3],L[4],L[5],L[6],L[7],L[8],L[9],"DB date is synthetic 2026-03-06 import","Live",L[10],L[11]])
ws.append(["TOTAL (12 locations)","4.5-4.8","",CAPTURED,ROUTED,PRIVATE,LOWSTAR,ROUTED/CAPTURED,PRIVATE/CAPTURED,"","captured = routed + held private (verified)","Live","901 = 606 (<=3) + 211 (4) + 84 (5)","Proven internal tool"])
style_header(ws,len(cols))
for r in range(2,ws.max_row+1):
    for cc in (3,4,5,6,7): ws.cell(r,cc).number_format="#,##0"
    ws.cell(r,8).number_format="0.0%"; ws.cell(r,9).number_format="0.0%"
    for cc in range(1,len(cols)+1): ws.cell(r,cc).border=border
for cc in range(1,len(cols)+1):
    ws.cell(ws.max_row,cc).fill=tot_fill; ws.cell(ws.max_row,cc).font=Font(bold=True)
autofit(ws)
ws.column_dimensions["M"].width=46; ws.column_dimensions["N"].width=44

# Tab B claim audit
ws2=wb.create_sheet("Claim Audit")
c2=["Claim","Use in board deck?","Verified fact","Confidence","Risk","Safe wording"]
ws2.append(c2)
for r in CLAIMS: ws2.append(r)
style_header(ws2,len(c2))
for r in range(2,ws2.max_row+1):
    for cc in range(1,len(c2)+1): ws2.cell(r,cc).border=border; ws2.cell(r,cc).alignment=Alignment(wrap_text=True,vertical="top")
autofit(ws2)
ws2.column_dimensions["A"].width=42; ws2.column_dimensions["C"].width=40; ws2.column_dimensions["F"].width=46

# Tab C open questions
ws3=wb.create_sheet("Open Questions")
c3=["Open question","What to confirm","Why it's open","Owner"]
ws3.append(c3)
for q in OPEN_Q: ws3.append(q)
style_header(ws3,len(c3))
for r in range(2,ws3.max_row+1):
    for cc in range(1,len(c3)+1): ws3.cell(r,cc).border=border; ws3.cell(r,cc).alignment=Alignment(wrap_text=True,vertical="top")
autofit(ws3)
for col,wd in [("A",30),("B",40),("C",34),("D",26)]: ws3.column_dimensions[col].width=wd

xlsx_path=os.path.join(OUT,"RateTap_Location_Appendix_2026-06-16.xlsx"); wb.save(xlsx_path)
print("XLSX:", xlsx_path, "tabs:", wb.sheetnames)

# ════════════════════════════════════════════════════════════════════════
# 4. Memo reply TXT
# ════════════════════════════════════════════════════════════════════════
memo=f"""=== WHATSAPP VERSION ===

Memo, pulled the current numbers straight from production (not the older case-studies page).

Quick note on the "2 missing locations": all 12 locations ARE in the case-studies summary table. Only 3 have detailed write-up cards (Angelopolis, Juarez, Regio Norte). So the gap is detailed cards (9 missing), not missing locations.

Board-safe story:
- RateTap is live across all 12 locations
- {fmt(CAPTURED)} reviews captured
- {fmt(ROUTED)} guests routed toward Google
- {fmt(LOWSTAR)} low ratings (3 stars or below) routed privately
- Ratings 4.5-4.8

Important: present it as a working internal tool, not external SaaS revenue - there are no active external paying customers today, and our one external pilot (El Braserio) didn't convert. Also we should NOT claim RateTap raised Google ratings (early growth was a Google place-ID artifact).

Full board package is ready. One open item before Tuesday: confirm in Stripe whether El Braserio's setup fee was actually collected - I can't see that from our database.


=== EMAIL VERSION ===

Hi Memo,

I pulled the current numbers directly from the production database, so these are live as of June 6, 2026.

On the "2 missing locations": all 12 locations already appear in the case-studies summary table. Only 3 have detailed write-up cards (Estancia Angelopolis, Estancia Juarez, Regio Norte). So there aren't 2 missing locations - there are 9 locations without a detailed card. Want me to build cards for those 9, or did you mean something else?

The board-safe story:
- RateTap is live across all 12 Grupo Estancia locations and used daily (856 reviews in the last 7 days).
- It has captured {fmt(CAPTURED)} reviews, routed {fmt(ROUTED)} guests toward Google, and routed {fmt(LOWSTAR)} low ratings (3 stars or below) to a private feedback path.
- Group Google ratings are 4.5-4.8 stars.

Please keep these caveats visible (they matter in the room):
- It is operationally working, but not an external SaaS business yet - no active external paying customers, and our one external pilot (El Braserio) did not convert.
- We should not claim RateTap increased Google ratings - the earlier growth figures were affected by a Google place-ID artifact.
- The big Google review counts on our profiles are lifetime totals, not RateTap's output. RateTap's defensible contribution is the {fmt(ROUTED)} routed.
- "Held private" is not always negative - only {fmt(LOWSTAR)} of the {fmt(PRIVATE)} are low ratings.

One thing I need a human to confirm before Tuesday: whether El Braserio's one-time setup fee was actually collected. There's no payment record in our database, so it can only be confirmed in the Stripe dashboard. I don't want to state a revenue number we can't back up.

I've prepared a full presentation-ready package (one-page brief, slide deck, board packet PDF, Excel appendix, and speaker notes). It's ready for your review.

Best,
Lawrence
"""
memo_txt=os.path.join(OUT,"Memo_Reply_Ready_To_Send.txt")
with open(memo_txt,"w") as f: f.write(memo)
print("TXT:", memo_txt)

print("\nALL DELIVERABLES BUILT OK")
